"""
PolyDoc AI - Multi-format Document Processor
Handles PDF, DOCX, PPTX, and image files with layout preservation
"""

import os
import asyncio
from typing import Dict, List, Any, Optional, Tuple
from dataclasses import dataclass
from pathlib import Path
import logging

# Document processing imports
import PyPDF2
from docx import Document as DocxDocument
from pptx import Presentation
from PIL import Image
import cv2
import numpy as np

# OCR imports
import pytesseract
import easyocr

# Layout analysis (optional)
try:
    import layoutparser as lp
    LAYOUT_PARSER_AVAILABLE = True
except ImportError:
    lp = None
    LAYOUT_PARSER_AVAILABLE = False

@dataclass
class DocumentElement:
    """Represents a document element with layout information"""
    text: str
    page_number: int
    element_type: str  # 'paragraph', 'heading', 'table', 'image', 'handwriting'
    bbox: Tuple[float, float, float, float]  # (x1, y1, x2, y2)
    confidence: float
    language: Optional[str] = None
    font_info: Optional[Dict] = None

@dataclass
class ProcessedDocument:
    """Container for processed document information"""
    filename: str
    total_pages: int
    elements: List[DocumentElement]
    summary: Optional[str] = None
    metadata: Dict[str, Any] = None

class DocumentProcessor:
    """Main document processor supporting multiple formats"""
    
    def __init__(self):
        self.logger = logging.getLogger(__name__)
        
        # Initialize EasyOCR for multilingual support
        # Using a minimal set to avoid compatibility issues
        self.ocr_reader = easyocr.Reader(['en'])
        
        # Initialize layout model if layoutparser is available
        self.layout_model = None
        if LAYOUT_PARSER_AVAILABLE:
            try:
                self.layout_model = lp.Detectron2LayoutModel(
                    'lp://PubLayNet/faster_rcnn_R_50_FPN_3x/config',
                    extra_config=["MODEL.ROI_HEADS.SCORE_THRESH_TEST", 0.8],
                    label_map={0: "Text", 1: "Title", 2: "List", 3: "Table", 4: "Figure"}
                )
            except Exception as e:
                self.logger.warning(f"Layout model not available: {e}")
                self.layout_model = None
        else:
            self.logger.info("Layout parser not available, skipping layout model initialization")
        
        # Supported file types
        self.supported_formats = {'.pdf', '.docx', '.pptx', '.png', '.jpg', '.jpeg', '.tiff', '.bmp'}
    
    def estimate_processing_time(self, file_path: str) -> dict:
        """Estimate processing time based on file size and type"""
        try:
            file_path = Path(file_path)
            file_size_mb = file_path.stat().st_size / (1024 * 1024)  # Size in MB
            file_ext = file_path.suffix.lower()
            
            # Base processing times (seconds per MB)
            time_estimates = {
                '.pdf': 3,      # PDF processing is moderate
                '.docx': 1,     # DOCX is fastest
                '.pptx': 2,     # PPTX is moderate
                '.png': 8,      # Images need OCR - slower
                '.jpg': 8,
                '.jpeg': 8,
                '.tiff': 10,    # TIFF can be large
                '.bmp': 12      # BMP is uncompressed - slowest
            }
            
            base_time = time_estimates.get(file_ext, 8)  # Default to 8s/MB
            estimated_seconds = max(5, int(file_size_mb * base_time))  # Minimum 5 seconds
            
            return {
                'estimated_seconds': estimated_seconds,
                'estimated_minutes': round(estimated_seconds / 60, 1),
                'file_size_mb': round(file_size_mb, 2),
                'complexity': 'High' if estimated_seconds > 60 else 'Medium' if estimated_seconds > 20 else 'Low'
            }
        except Exception as e:
            self.logger.error(f"Error estimating time: {e}")
            return {
                'estimated_seconds': 30,
                'estimated_minutes': 0.5,
                'file_size_mb': 0,
                'complexity': 'Unknown'
            }
    
    async def process_document(self, file_path: str) -> ProcessedDocument:
        """Main method to process any supported document format"""
        try:
            file_path = Path(file_path)
            
            if not file_path.exists():
                raise FileNotFoundError(f"File not found: {file_path}")
            
            if file_path.suffix.lower() not in self.supported_formats:
                raise ValueError(f"Unsupported format: {file_path.suffix}")
            
            self.logger.info(f"Processing document: {file_path.name}")
            
            # Route to appropriate processor
            if file_path.suffix.lower() == '.pdf':
                return await self._process_pdf(file_path)
            elif file_path.suffix.lower() == '.docx':
                return await self._process_docx(file_path)
            elif file_path.suffix.lower() == '.pptx':
                return await self._process_pptx(file_path)
            else:  # Image formats
                return await self._process_image(file_path)
        
        except Exception as e:
            self.logger.error(f"Error processing document: {e}")
            raise
    
    async def _process_pdf(self, file_path: Path) -> ProcessedDocument:
        """Process PDF files with layout preservation"""
        elements = []
        
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            total_pages = len(pdf_reader.pages)
            
            for page_num, page in enumerate(pdf_reader.pages, 1):
                # Extract text
                text_content = page.extract_text()
                
                if text_content.strip():
                    # For PDFs, we'll treat the extracted text as paragraphs
                    paragraphs = [p.strip() for p in text_content.split('\n\n') if p.strip()]
                    
                    for i, paragraph in enumerate(paragraphs):
                        element = DocumentElement(
                            text=paragraph,
                            page_number=page_num,
                            element_type='paragraph',
                            bbox=(0, i*50, 500, (i+1)*50),  # Approximate positioning
                            confidence=0.9,
                            language=self._detect_language(paragraph)
                        )
                        elements.append(element)
                
                # Convert PDF page to image for layout analysis if needed
                # This would require pdf2image library for more advanced layout detection
        
        return ProcessedDocument(
            filename=file_path.name,
            total_pages=total_pages,
            elements=elements,
            metadata={'file_type': 'pdf', 'size': file_path.stat().st_size}
        )
    
    async def _process_docx(self, file_path: Path) -> ProcessedDocument:
        """Process DOCX files with layout preservation"""
        doc = DocxDocument(file_path)
        elements = []
        
        for para_idx, paragraph in enumerate(doc.paragraphs):
            if paragraph.text.strip():
                # Determine element type based on style
                element_type = 'heading' if paragraph.style.name.startswith('Heading') else 'paragraph'
                
                element = DocumentElement(
                    text=paragraph.text,
                    page_number=1,  # DOCX doesn't have clear page boundaries
                    element_type=element_type,
                    bbox=(0, para_idx*30, 500, (para_idx+1)*30),
                    confidence=1.0,
                    language=self._detect_language(paragraph.text),
                    font_info={'style': paragraph.style.name}
                )
                elements.append(element)
        
        # Process tables
        for table_idx, table in enumerate(doc.tables):
            table_text = []
            for row in table.rows:
                row_text = [cell.text for cell in row.cells]
                table_text.append(' | '.join(row_text))
            
            if table_text:
                element = DocumentElement(
                    text='\n'.join(table_text),
                    page_number=1,
                    element_type='table',
                    bbox=(0, 1000+table_idx*100, 500, 1000+(table_idx+1)*100),
                    confidence=1.0,
                    language=self._detect_language(' '.join(table_text))
                )
                elements.append(element)
        
        return ProcessedDocument(
            filename=file_path.name,
            total_pages=1,
            elements=elements,
            metadata={'file_type': 'docx', 'size': file_path.stat().st_size}
        )
    
    async def _process_pptx(self, file_path: Path) -> ProcessedDocument:
        """Process PowerPoint files"""
        prs = Presentation(file_path)
        elements = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    # Determine element type
                    element_type = 'heading' if shape.text.isupper() or len(shape.text) < 50 else 'paragraph'
                    
                    element = DocumentElement(
                        text=shape.text,
                        page_number=slide_num,
                        element_type=element_type,
                        bbox=(shape.left, shape.top, shape.left + shape.width, shape.top + shape.height),
                        confidence=1.0,
                        language=self._detect_language(shape.text)
                    )
                    elements.append(element)
        
        return ProcessedDocument(
            filename=file_path.name,
            total_pages=len(prs.slides),
            elements=elements,
            metadata={'file_type': 'pptx', 'size': file_path.stat().st_size}
        )
    
    async def _process_image(self, file_path: Path) -> ProcessedDocument:
        """Process image files with OCR and layout analysis"""
        # Load image
        image = cv2.imread(str(file_path))
        if image is None:
            raise ValueError(f"Could not load image: {file_path}")
        
        elements = []
        
        # Perform OCR with EasyOCR
        ocr_results = self.ocr_reader.readtext(str(file_path))
        
        for result in ocr_results:
            bbox_coords = result[0]
            text = result[1]
            confidence = result[2]
            
            if confidence > 0.5 and text.strip():
                # Convert bbox format
                x_coords = [coord[0] for coord in bbox_coords]
                y_coords = [coord[1] for coord in bbox_coords]
                bbox = (min(x_coords), min(y_coords), max(x_coords), max(y_coords))
                
                # Determine if it's handwriting based on confidence and other factors
                element_type = 'handwriting' if confidence < 0.8 else 'text'
                
                element = DocumentElement(
                    text=text,
                    page_number=1,
                    element_type=element_type,
                    bbox=bbox,
                    confidence=confidence,
                    language=self._detect_language(text)
                )
                elements.append(element)
        
        # Perform layout analysis if model is available
        if self.layout_model:
            try:
                layout = self.layout_model.detect(image)
                # This would enhance the element classification
                # Implementation depends on specific layout model output
            except Exception as e:
                self.logger.warning(f"Layout analysis failed: {e}")
        
        return ProcessedDocument(
            filename=file_path.name,
            total_pages=1,
            elements=elements,
            metadata={'file_type': 'image', 'size': file_path.stat().st_size}
        )
    
    def _detect_language(self, text: str) -> str:
        """Simple language detection based on character analysis"""
        if not text:
            return 'unknown'
        
        # Count different script characters
        latin_chars = sum(1 for c in text if c.isascii() and c.isalpha())
        arabic_chars = sum(1 for c in text if '\u0600' <= c <= '\u06FF')
        hindi_chars = sum(1 for c in text if '\u0900' <= c <= '\u097F')
        chinese_chars = sum(1 for c in text if '\u4e00' <= c <= '\u9fff')
        
        total_chars = latin_chars + arabic_chars + hindi_chars + chinese_chars
        
        if total_chars == 0:
            return 'unknown'
        
        # Determine dominant script
        if arabic_chars / total_chars > 0.3:
            return 'ar'
        elif hindi_chars / total_chars > 0.3:
            return 'hi'
        elif chinese_chars / total_chars > 0.3:
            return 'zh'
        else:
            return 'en'
    
    def get_document_stats(self, document: ProcessedDocument) -> Dict[str, Any]:
        """Generate statistics about the processed document"""
        stats = {
            'total_elements': len(document.elements),
            'total_pages': document.total_pages,
            'element_types': {},
            'languages': {},
            'avg_confidence': 0,
            'total_text_length': 0
        }
        
        for element in document.elements:
            # Count element types
            stats['element_types'][element.element_type] = \
                stats['element_types'].get(element.element_type, 0) + 1
            
            # Count languages
            if element.language:
                stats['languages'][element.language] = \
                    stats['languages'].get(element.language, 0) + 1
            
            # Calculate averages
            stats['avg_confidence'] += element.confidence
            stats['total_text_length'] += len(element.text)
        
        if document.elements:
            stats['avg_confidence'] /= len(document.elements)
        
        return stats
